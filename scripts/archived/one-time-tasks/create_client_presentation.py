#!/usr/bin/env python3
"""
ZTL数智化作战中心 - 客户商务展示PPT生成脚本
设计风格：Cyberpunk风格 + 商务专业性
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# 定义Cyberpunk配色方案（商务化调整）
colors = {
    'neon_cyan': RGBColor(0, 255, 255),
    'neon_pink': RGBColor(255, 0, 255),
    'neon_purple': RGBColor(176, 0, 255),
    'neon_yellow': RGBColor(255, 255, 0),
    'neon_green': RGBColor(0, 255, 0),
    'neon_orange': RGBColor(255, 136, 0),
    'dark_bg': RGBColor(10, 10, 15),
    'secondary_bg': RGBColor(26, 10, 46),
    'text_light': RGBColor(168, 168, 255),
    'white': RGBColor(255, 255, 255),
    'gray': RGBColor(136, 136, 204)
}

def create_presentation():
    """创建演示文稿基础对象"""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    return prs

def set_dark_background(slide):
    """设置深色背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['dark_bg']

def add_title_text(slide, x, y, width, height, text, font_size=48, color=None, bold=True, alignment=PP_ALIGN.CENTER):
    """添加标题文本框"""
    if color is None:
        color = colors['neon_cyan']

    text_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    text_frame = text_box.text_frame
    text_frame.text = text
    para = text_frame.paragraphs[0]
    para.font.name = 'Microsoft YaHei UI'
    para.font.size = Pt(font_size)
    para.font.bold = bold
    para.font.color.rgb = color
    para.alignment = alignment
    return text_frame

def add_body_text(slide, x, y, width, height, text, font_size=18, color=None, alignment=PP_ALIGN.LEFT):
    """添加正文文本框"""
    if color is None:
        color = colors['text_light']

    text_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    text_frame = text_box.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    para = text_frame.paragraphs[0]
    para.font.name = 'Microsoft YaHei UI'
    para.font.size = Pt(font_size)
    para.font.color.rgb = color
    para.alignment = alignment
    para.line_spacing = 1.5
    return text_frame

def add_bullet_points(slide, x, y, width, height, points, font_size=20, color=None):
    """添加项目符号列表"""
    if color is None:
        color = colors['text_light']

    text_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    for i, point in enumerate(points):
        if i > 0:
            para = text_frame.add_paragraph()
        else:
            para = text_frame.paragraphs[0]

        para.text = point
        para.font.name = 'Microsoft YaHei UI'
        para.font.size = Pt(font_size)
        para.font.color.rgb = color
        para.level = 0
        para.line_spacing = 1.5

    return text_frame

def create_slide_1_cover(prs):
    """第1页：封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)

    # 主标题
    add_title_text(slide, 1, 3, 14, 1.5, 'ZTL数智化作战中心', font_size=72, color=colors['neon_cyan'])

    # 副标题
    add_title_text(slide, 1, 5, 14, 1, '下一代AI智能体编排平台 · 赋能餐饮行业数智化转型',
                   font_size=28, color=colors['neon_pink'])

    # 版本信息
    add_body_text(slide, 5.5, 7.2, 5, 0.5, 'VERSION 2.1 | POWERED BY CLAUDE SONNET 4.5',
                  font_size=14, color=colors['gray'], alignment=PP_ALIGN.CENTER)

    print('✅ 第1页：封面 - 已创建')

def create_slide_2_problem(prs):
    """第2页：传统困境"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)

    # 标题
    add_title_text(slide, 1, 0.5, 14, 1, '传统餐饮管理的痛点', font_size=48, color=colors['neon_pink'])

    # 左侧：问题列表
    problems = [
        '⏰ 人工操作效率低下：8小时工作需要3人协作',
        '💰 人力成本高昂：专业人员招聘难、培养慢',
        '❌ 错误率居高不下：人工处理错误率达15%',
        '📉 经验难以复用：个人能力难以标准化传承',
        '🔄 流程管理混乱：缺乏系统化的协作机制'
    ]
    add_bullet_points(slide, 1.5, 2, 12, 5, problems, font_size=24, color=colors['text_light'])

    # 底部总结
    add_title_text(slide, 1, 7.5, 14, 0.8, '传统方式已无法满足现代餐饮企业的快速发展需求',
                   font_size=24, color=colors['neon_orange'], bold=False)

    print('✅ 第2页：传统困境 - 已创建')

def create_slide_3_solution(prs):
    """第3页：解决方案"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)

    # 标题
    add_title_text(slide, 1, 0.5, 14, 1, 'ZTL数智化作战中心：AI驱动的智能解决方案',
                   font_size=44, color=colors['neon_cyan'])

    # 核心价值主张
    value_props = [
        '🤖 60个专业AI智能体协同工作',
        '⚡ 时间效率提升70% - 从8小时降至2.4小时',
        '💎 成本降低67% - 从3人协作到1人监督',
        '✅ 错误率下降80% - 从15%降至3%',
        '♻️ 可复用性提升233% - 模块化设计，持续优化'
    ]
    add_bullet_points(slide, 1.5, 2.5, 13, 4.5, value_props, font_size=26, color=colors['neon_green'])

    # 底部标语
    add_title_text(slide, 1, 7.5, 14, 0.8, '让AI智能体成为您的数字化员工',
                   font_size=28, color=colors['neon_pink'], bold=True)

    print('✅ 第3页：解决方案 - 已创建')

def create_slide_4_capabilities(prs):
    """第4页：核心能力"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)

    # 标题
    add_title_text(slide, 1, 0.4, 14, 0.8, '7大业务组 · 全场景覆盖', font_size=48, color=colors['neon_cyan'])

    # 业务组卡片 - 3行，第一行3个，第二行3个，第三行1个
    business_groups = [
        ('战略组', 'G系列 · 9个智能体\n经营分析·产品力·选址评估', colors['neon_pink']),
        ('创意组', 'X系列 · 9个智能体\n广告策划·内容创作·视觉设计', colors['neon_yellow']),
        ('情报组', 'E系列 · 8个智能体\n数据采集·深度分析·云端存储', colors['neon_cyan']),
        ('行政组', 'R系列 · 8个智能体\n财务人事·协同管理·文档处理', colors['neon_green']),
        ('中台组', 'M系列 · 7个智能体\n美团运营·供应链·网页自动化', colors['neon_orange']),
        ('筹建组', 'Z系列 · 6个智能体\nBIM建模·空间设计·3D渲染', RGBColor(0, 255, 136)),
        ('开发组', 'F系列 · 11个智能体\n全栈开发·技术架构·云端部署', colors['neon_purple'])
    ]

    positions = [
        # 第一行
        (1.5, 1.8, 4, 1.8),
        (6, 1.8, 4, 1.8),
        (10.5, 1.8, 4, 1.8),
        # 第二行
        (1.5, 4, 4, 1.8),
        (6, 4, 4, 1.8),
        (10.5, 4, 4, 1.8),
        # 第三行（居中）
        (6, 6.2, 4, 1.8)
    ]

    for i, ((title, desc, color), (x, y, w, h)) in enumerate(zip(business_groups, positions)):
        # 标题
        add_title_text(slide, x, y, w, 0.5, title, font_size=24, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        # 描述
        add_body_text(slide, x, y + 0.5, w, h - 0.5, desc, font_size=14, color=colors['text_light'], alignment=PP_ALIGN.CENTER)

    print('✅ 第4页：核心能力 - 已创建')

def create_slide_5_architecture(prs):
    """第5页：技术架构（简化）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)

    # 标题
    add_title_text(slide, 1, 0.4, 14, 0.8, '三层智能架构', font_size=48, color=colors['neon_cyan'])

    # 架构层级
    layers = [
        ('知识层',
         '60个专业智能体 + 34个技能包\n提供决策框架和工具能力',
         colors['neon_cyan'],
         1.5, 1.8),
        ('决策层',
         'Claude Sonnet 4.5运行时\n智能理解需求·动态编排执行',
         colors['neon_pink'],
         1.5, 3.8),
        ('执行层',
         '8个MCP服务器 + 多类型技能包\n真实执行工具·持续输出结果',
         colors['neon_purple'],
         1.5, 5.8)
    ]

    for i, (title, desc, color, x, y) in enumerate(layers):
        # 层级标题
        add_title_text(slide, x, y, 12, 0.5, f'第{i+1}层：{title}', font_size=28, color=color, alignment=PP_ALIGN.LEFT)
        # 层级描述
        add_body_text(slide, x + 0.3, y + 0.5, 12, 1, desc, font_size=18, color=colors['text_light'])

        # 箭头（除最后一层）
        if i < len(layers) - 1:
            add_title_text(slide, 7, y + 1.6, 1, 0.4, '↓', font_size=36, color=colors['neon_pink'], alignment=PP_ALIGN.CENTER)

    print('✅ 第5页：技术架构 - 已创建')

def create_slide_6_case_studies(prs):
    """第6页：实战案例"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)

    # 标题
    add_title_text(slide, 1, 0.4, 14, 0.8, '实战案例：效率倍增的成功实践', font_size=44, color=colors['neon_cyan'])

    # 案例卡片
    cases = [
        ('新品上市Campaign',
         '传统方式：8小时·3人协作\nAI智能体：2.4小时·1人监督\n产出：7类营销物料一次性交付',
         colors['neon_pink']),
        ('竞品情报分析',
         '传统方式：16小时·多工具切换\nAI智能体：4小时·自动化采集\n覆盖：20+数据源·50+分析指标',
         colors['neon_cyan']),
        ('门店筹建全流程',
         '传统方式：15天·多方协调\nAI智能体：5天·流程自动化\n交付：CAD图纸+BIM模型+VR漫游',
         colors['neon_purple'])
    ]

    positions = [(1.5, 2), (6, 2), (10.5, 2)]

    for (title, desc, color), (x, y) in zip(cases, positions):
        # 案例标题
        add_title_text(slide, x, y, 4, 0.5, title, font_size=20, color=color, alignment=PP_ALIGN.CENTER)
        # 案例描述
        add_body_text(slide, x, y + 0.5, 4, 4, desc, font_size=14, color=colors['text_light'], alignment=PP_ALIGN.LEFT)

    print('✅ 第6页：实战案例 - 已创建')

def create_slide_7_efficiency(prs):
    """第7页：效率对比"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)

    # 标题
    add_title_text(slide, 1, 0.4, 14, 0.8, '效率对比：传统 VS AI智能体', font_size=48, color=colors['neon_cyan'])

    # 左侧：传统方式
    add_title_text(slide, 2, 1.8, 5, 0.6, '传统方式', font_size=32, color=colors['gray'], alignment=PP_ALIGN.CENTER)
    traditional_metrics = [
        '⏱️ 任务完成时间：8小时（100%）',
        '👥 人力成本：3人协作（100%）',
        '❌ 错误率：15%',
        '♻️ 可复用性：30%（低）'
    ]
    add_bullet_points(slide, 2, 2.8, 5, 4, traditional_metrics, font_size=18, color=colors['text_light'])

    # 中间：VS标识
    add_title_text(slide, 7, 4, 2, 0.8, 'VS', font_size=48, color=colors['neon_pink'], bold=True)

    # 右侧：AI智能体
    add_title_text(slide, 9, 1.8, 5, 0.6, 'AI智能体协作', font_size=32, color=colors['neon_cyan'], alignment=PP_ALIGN.CENTER)
    ai_metrics = [
        '⚡ 任务完成时间：2.4小时（↓70%）',
        '🤖 人力成本：1人监督（↓67%）',
        '✅ 错误率：3%（↓80%）',
        '🚀 可复用性：100%（↑233%）'
    ]
    add_bullet_points(slide, 9, 2.8, 5, 4, ai_metrics, font_size=18, color=colors['neon_green'])

    print('✅ 第7页：效率对比 - 已创建')

def create_slide_8_value(prs):
    """第8页：价值总结"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)

    # 标题
    add_title_text(slide, 1, 0.4, 14, 0.8, '为什么选择ZTL数智化作战中心？', font_size=44, color=colors['neon_cyan'])

    # ROI核心指标 - 2x3布局
    roi_metrics = [
        ('70%', '时间节省', colors['neon_cyan']),
        ('67%', '成本降低', colors['neon_pink']),
        ('80%', '错误率下降', colors['neon_purple']),
        ('233%', '复用性提升', colors['neon_green']),
        ('99.8%', '系统稳定性', colors['neon_orange']),
        ('3.2s', '平均响应', RGBColor(0, 255, 136))
    ]

    positions = [
        (2, 2, 3.5, 1.5),
        (6.25, 2, 3.5, 1.5),
        (10.5, 2, 3.5, 1.5),
        (2, 4.5, 3.5, 1.5),
        (6.25, 4.5, 3.5, 1.5),
        (10.5, 4.5, 3.5, 1.5)
    ]

    for (value, label, color), (x, y, w, h) in zip(roi_metrics, positions):
        # 数值
        add_title_text(slide, x, y, w, h * 0.6, value, font_size=42, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        # 标签
        add_body_text(slide, x, y + h * 0.6, w, h * 0.4, label, font_size=16, color=colors['text_light'], alignment=PP_ALIGN.CENTER)

    # 底部总结
    add_title_text(slide, 1, 7.2, 14, 0.8, '投资回报周期：3-6个月 | 长期ROI：300%+',
                   font_size=24, color=colors['neon_pink'], bold=True)

    print('✅ 第8页：价值总结 - 已创建')

def create_slide_9_contact(prs):
    """第9页：联系我们"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)

    # 主标题
    add_title_text(slide, 1, 2.5, 14, 1.2, '开启您的数智化转型之旅', font_size=54, color=colors['neon_cyan'])

    # 副标题
    add_title_text(slide, 1, 4, 14, 0.8, '立即联系我们，获取专属解决方案',
                   font_size=32, color=colors['neon_pink'], bold=False)

    # 联系信息
    contact_info = [
        '📧 邮箱：contact@ztl-ai.com',
        '📱 电话：400-XXX-XXXX',
        '🌐 官网：www.ztl-ai.com',
        '📍 地址：[您的公司地址]'
    ]
    add_bullet_points(slide, 4.5, 5.5, 7, 2, contact_info, font_size=20, color=colors['text_light'])

    # 底部
    add_body_text(slide, 5, 7.8, 6, 0.4, 'Powered by Claude Sonnet 4.5 | © 2025 ZTL数智化作战中心',
                  font_size=12, color=colors['gray'], alignment=PP_ALIGN.CENTER)

    print('✅ 第9页：联系我们 - 已创建')

def main():
    """主函数：生成完整PPT"""
    print('\n' + '='*60)
    print('ZTL数智化作战中心 - 客户商务展示PPT生成中...')
    print('='*60 + '\n')

    # 创建演示文稿
    prs = create_presentation()

    # 创建所有幻灯片
    create_slide_1_cover(prs)
    create_slide_2_problem(prs)
    create_slide_3_solution(prs)
    create_slide_4_capabilities(prs)
    create_slide_5_architecture(prs)
    create_slide_6_case_studies(prs)
    create_slide_7_efficiency(prs)
    create_slide_8_value(prs)
    create_slide_9_contact(prs)

    # 保存文件
    output_dir = 'output'
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, 'ZTL数智化作战中心-商务展示PPT.pptx')
    prs.save(output_path)

    print('\n' + '='*60)
    print(f'✅ PPT生成成功！')
    print(f'📁 文件位置：{output_path}')
    print(f'📊 包含9页精心设计的商务展示内容')
    print('='*60 + '\n')

if __name__ == '__main__':
    main()
