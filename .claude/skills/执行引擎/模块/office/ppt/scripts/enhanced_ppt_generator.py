"""
Enhanced PPT Generator with Professional 2025 Design Standards
集成渐变背景、阴影效果、现代配色和字体层级优化
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement


class EnhancedPPTGenerator:
    """增强版PPT生成器 - 2025专业设计标准"""

    # 活力营销风配色方案
    COLOR_PRIMARY = RGBColor(255, 107, 53)      # #FF6B35 活力橙
    COLOR_SECONDARY = RGBColor(0, 78, 137)       # #004E89 专业蓝
    COLOR_ACCENT = RGBColor(247, 179, 43)        # #F7B32B 金黄
    COLOR_DARK = RGBColor(26, 26, 26)            # #1A1A1A 深黑
    COLOR_LIGHT = RGBColor(255, 255, 255)        # #FFFFFF 纯白
    COLOR_GRAY = RGBColor(117, 117, 117)         # #757575 中灰
    COLOR_BG_LIGHT = RGBColor(250, 250, 250)     # #FAFAFA 浅灰背景

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)  # 16:9

    def _add_gradient_background(self, slide, color1, color2, angle=45):
        """添加渐变背景"""
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = angle

        # 设置渐变停止点
        fill.gradient_stops[0].color.rgb = color1
        fill.gradient_stops[1].color.rgb = color2

    def _add_shadow(self, shape, blur=6, distance=3, direction=315, transparency=30):
        """为形状添加柔和阴影"""
        shadow = shape.shadow
        shadow.inherit = False
        shadow.shadow_type = 'outer'
        shadow.blur_radius = Pt(blur)
        shadow.distance = Pt(distance)
        shadow.angle = direction
        shadow.transparency = transparency / 100.0
        return shadow

    def _set_text_format(self, text_frame, font_size, font_color, bold=False,
                         alignment=PP_ALIGN.LEFT, line_spacing=1.2):
        """设置文本格式"""
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = alignment
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(font_size * 0.3)
            paragraph.line_spacing = line_spacing

            # 如果段落没有runs,添加一个
            if not paragraph.runs:
                paragraph.text = paragraph.text  # 这会创建runs

            for run in paragraph.runs:
                run.font.name = 'Arial'  # 使用系统通用字体
                run.font.size = Pt(font_size)
                if font_color:
                    try:
                        run.font.color.rgb = font_color
                    except:
                        # 如果设置失败,使用默认颜色
                        pass
                run.font.bold = bold

    def add_title_slide(self, title, subtitle):
        """添加封面页 - 使用渐变背景"""
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)

        # 添加橙蓝渐变背景
        self._add_gradient_background(slide, self.COLOR_PRIMARY, self.COLOR_SECONDARY, 135)

        # 添加标题文本框
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(1.5)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 标题样式: 60pt 粗体 白色
        self._set_text_format(title_frame, 60, self.COLOR_LIGHT, bold=True,
                              alignment=PP_ALIGN.CENTER)

        # 添加副标题
        subtitle_box = slide.shapes.add_textbox(left, top + Inches(1.2), width, Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle

        # 副标题样式: 28pt 常规 白色
        self._set_text_format(subtitle_frame, 28, self.COLOR_LIGHT,
                              alignment=PP_ALIGN.CENTER)

        return slide

    def add_content_slide(self, title, content_list, use_icons=False):
        """添加内容页 - 带渐变标题区和阴影卡片"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 浅色背景
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.COLOR_BG_LIGHT

        # 添加标题区域 (带渐变)
        title_shape = slide.shapes.add_shape(
            1,  # 矩形
            Inches(0), Inches(0),
            Inches(10), Inches(0.8)
        )
        title_shape.fill.gradient()
        title_shape.fill.gradient_angle = 0
        title_shape.fill.gradient_stops[0].color.rgb = self.COLOR_PRIMARY
        title_shape.fill.gradient_stops[1].color.rgb = self.COLOR_SECONDARY
        title_shape.line.fill.background()

        # 标题文本
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.margin_left = Inches(0.3)
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._set_text_format(title_frame, 36, self.COLOR_LIGHT, bold=True)

        # 内容卡片 (带阴影)
        content_shape = slide.shapes.add_shape(
            1,  # 矩形
            Inches(0.5), Inches(1.2),
            Inches(9), Inches(3.9)
        )
        content_shape.fill.solid()
        content_shape.fill.fore_color.rgb = self.COLOR_LIGHT
        content_shape.line.fill.background()

        # 添加柔和阴影
        self._add_shadow(content_shape, blur=10, distance=5, transparency=20)

        # 内容文本
        content_frame = content_shape.text_frame
        content_frame.margin_left = Inches(0.4)
        content_frame.margin_right = Inches(0.4)
        content_frame.margin_top = Inches(0.3)
        content_frame.word_wrap = True

        for item in content_list:
            p = content_frame.add_paragraph()
            p.level = 0
            p.space_before = Pt(8)
            p.space_after = Pt(8)
            p.line_spacing = 1.3

            # 添加run并设置格式
            run = p.add_run()
            run.text = item
            font = run.font
            font.name = 'Arial'
            font.size = Pt(20)
            # 安全设置颜色 - 必须在文本赋值之后
            try:
                font.color.rgb = self.COLOR_DARK
            except (AttributeError, TypeError):
                pass  # 如果颜色设置失败，使用默认颜色

        return slide

    def add_two_column_slide(self, title, left_content, right_content):
        """添加双栏对比页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.COLOR_BG_LIGHT

        # 标题区
        title_shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(10), Inches(0.8)
        )
        title_shape.fill.gradient()
        title_shape.fill.gradient_angle = 0
        title_shape.fill.gradient_stops[0].color.rgb = self.COLOR_SECONDARY
        title_shape.fill.gradient_stops[1].color.rgb = self.COLOR_PRIMARY
        title_shape.line.fill.background()

        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.margin_left = Inches(0.3)
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._set_text_format(title_frame, 36, self.COLOR_LIGHT, bold=True)

        # 左栏卡片
        left_shape = slide.shapes.add_shape(
            1, Inches(0.5), Inches(1.2), Inches(4.3), Inches(3.9)
        )
        left_shape.fill.solid()
        left_shape.fill.fore_color.rgb = self.COLOR_LIGHT
        left_shape.line.fill.background()
        self._add_shadow(left_shape, blur=10, distance=5, transparency=20)

        left_frame = left_shape.text_frame
        left_frame.margin_left = Inches(0.3)
        left_frame.margin_top = Inches(0.2)
        left_frame.word_wrap = True

        for item in left_content:
            p = left_frame.add_paragraph()
            p.space_before = Pt(6)
            p.space_after = Pt(6)
            p.line_spacing = 1.2
            run = p.add_run()
            run.text = item
            font = run.font
            font.name = 'Arial'
            font.size = Pt(18)
            try:
                font.color.rgb = self.COLOR_DARK
            except (AttributeError, TypeError):
                pass

        # 右栏卡片
        right_shape = slide.shapes.add_shape(
            1, Inches(5.2), Inches(1.2), Inches(4.3), Inches(3.9)
        )
        right_shape.fill.solid()
        right_shape.fill.fore_color.rgb = self.COLOR_LIGHT
        right_shape.line.fill.background()
        self._add_shadow(right_shape, blur=10, distance=5, transparency=20)

        right_frame = right_shape.text_frame
        right_frame.margin_left = Inches(0.3)
        right_frame.margin_top = Inches(0.2)
        right_frame.word_wrap = True

        for item in right_content:
            p = right_frame.add_paragraph()
            p.space_before = Pt(6)
            p.space_after = Pt(6)
            p.line_spacing = 1.2
            run = p.add_run()
            run.text = item
            font = run.font
            font.name = 'Arial'
            font.size = Pt(18)
            try:
                font.color.rgb = self.COLOR_DARK
            except (AttributeError, TypeError):
                pass

        return slide

    def add_table_slide(self, title, headers, rows, highlight_header=True):
        """添加表格页 - 使用渐变表头"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.COLOR_BG_LIGHT

        # 标题
        title_shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(10), Inches(0.8)
        )
        title_shape.fill.gradient()
        title_shape.fill.gradient_angle = 0
        title_shape.fill.gradient_stops[0].color.rgb = self.COLOR_PRIMARY
        title_shape.fill.gradient_stops[1].color.rgb = self.COLOR_ACCENT
        title_shape.line.fill.background()

        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.margin_left = Inches(0.3)
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._set_text_format(title_frame, 36, self.COLOR_LIGHT, bold=True)

        # 表格
        rows_count = len(rows) + 1  # +1 for header
        cols_count = len(headers)

        left = Inches(0.8)
        top = Inches(1.3)
        width = Inches(8.4)
        height = Inches(3.5)

        table = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table

        # 表头
        for col_idx, header_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.COLOR_SECONDARY

            text_frame = cell.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(header_text)
            font = run.font
            font.name = 'Arial'
            font.size = Pt(16)
            font.bold = True
            try:
                font.color.rgb = self.COLOR_LIGHT
            except (AttributeError, TypeError):
                pass

        # 数据行
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)

                # 斑马纹
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(245, 245, 245)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.COLOR_LIGHT

                text_frame = cell.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = str(cell_text)
                font = run.font
                font.name = 'Arial'
                font.size = Pt(14)
                try:
                    font.color.rgb = self.COLOR_DARK
                except (AttributeError, TypeError):
                    pass

        return slide

    def add_image_slide(self, title, image_path, caption=""):
        """添加图片页 - 居中展示"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.COLOR_BG_LIGHT

        # 标题
        title_shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(10), Inches(0.7)
        )
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = self.COLOR_DARK
        title_shape.line.fill.background()

        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.margin_left = Inches(0.3)
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._set_text_format(title_frame, 32, self.COLOR_LIGHT, bold=True)

        # 图片
        try:
            left = Inches(1.5)
            top = Inches(1.2)
            width = Inches(7)

            pic = slide.shapes.add_picture(image_path, left, top, width=width)

            # 添加图片阴影
            self._add_shadow(pic, blur=15, distance=8, transparency=25)

            # 图片说明
            if caption:
                caption_box = slide.shapes.add_textbox(
                    Inches(1.5), Inches(4.8), Inches(7), Inches(0.4)
                )
                caption_frame = caption_box.text_frame
                caption_frame.text = caption
                self._set_text_format(caption_frame, 14, self.COLOR_GRAY,
                                     alignment=PP_ALIGN.CENTER)
        except Exception as e:
            # 如果图片加载失败,添加占位文本
            error_box = slide.shapes.add_textbox(
                Inches(2), Inches(2.5), Inches(6), Inches(1)
            )
            error_frame = error_box.text_frame
            error_frame.text = f"图片加载失败: {image_path}"
            self._set_text_format(error_frame, 18, self.COLOR_GRAY,
                                 alignment=PP_ALIGN.CENTER)

        return slide

    def save(self, filename):
        """保存PPT文件"""
        self.prs.save(filename)
        return filename


# 测试代码
if __name__ == "__main__":
    generator = EnhancedPPTGenerator()

    # 测试封面
    generator.add_title_slide(
        "云南过桥米线品牌营销策划方案",
        '"一碗过桥,三餐云南" | 2025年10月'
    )

    # 测试内容页
    generator.add_content_slide(
        "核心营销策略",
        [
            "场景营销 - 早中晚三餐场景创新",
            "文化体验 - 云南非遗故事传播",
            "社交传播 - UGC内容共创机制"
        ]
    )

    # 测试双栏页
    generator.add_two_column_slide(
        "目标客群画像",
        ["核心人群: 25-40岁都市白领", "追求品质生活", "社交媒体活跃用户"],
        ["次级人群: 18-25岁Z世代", "喜欢新鲜体验", "拍照分享达人"]
    )

    # 保存
    generator.save("test_enhanced_ppt.pptx")
    print("✅ 增强版PPT生成器测试完成!")
