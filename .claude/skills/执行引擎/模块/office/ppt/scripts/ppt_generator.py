"""
PowerPoint Generator - Direct Generation Method

This module provides a fluent API for creating PowerPoint presentations
programmatically using python-pptx.
"""

import logging
from pathlib import Path
from typing import List, Dict, Any, Optional
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class PPTGenerator:
    """
    PowerPoint presentation builder with fluent API.

    Provides methods for creating various types of slides including:
    - Title slides
    - Content slides (bullet points)
    - Two-column slides
    - Table slides
    - Image slides
    - Chart slides
    """

    def __init__(self, template_path: Optional[str] = None):
        """
        Initialize the PPT generator.

        Args:
            template_path: Path to PowerPoint template file (.pptx).
                          If None, creates a blank presentation.
        """
        if template_path:
            self.prs = Presentation(template_path)
            logger.info(f"Loaded template: {template_path}")
        else:
            self.prs = Presentation()
            logger.info("Created new blank presentation")

        self.slide_layouts = self.prs.slide_layouts
        self.slides_created = 0

    def add_title_slide(
        self,
        title: str,
        subtitle: Optional[str] = None
    ) -> "PPTGenerator":
        """
        Add a title slide to the presentation.

        Args:
            title: Main title text
            subtitle: Subtitle text (optional)

        Returns:
            Self for method chaining
        """
        slide = self.prs.slides.add_slide(self.slide_layouts[0])

        # Set title
        slide.shapes.title.text = title
        logger.info(f"Added title slide: {title}")

        # Set subtitle if provided
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
            logger.info(f"Added subtitle: {subtitle}")

        self.slides_created += 1
        return self

    def add_content_slide(
        self,
        title: str,
        content: List[str],
        font_size: int = 18
    ) -> "PPTGenerator":
        """
        Add a content slide with bullet points.

        Args:
            title: Slide title
            content: List of bullet points
            font_size: Font size for content (default: 18)

        Returns:
            Self for method chaining
        """
        slide = self.prs.slides.add_slide(self.slide_layouts[1])
        slide.shapes.title.text = title

        # Add bullet points
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()

        for item in content:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(font_size)

        logger.info(f"Added content slide: {title} ({len(content)} points)")
        self.slides_created += 1
        return self

    def add_two_column_slide(
        self,
        title: str,
        left_content: List[str],
        right_content: List[str],
        font_size: int = 16
    ) -> "PPTGenerator":
        """
        Add a two-column content slide.

        Args:
            title: Slide title
            left_content: List of items for left column
            right_content: List of items for right column
            font_size: Font size for content (default: 16)

        Returns:
            Self for method chaining
        """
        # Try to use two-column layout (usually layout index 3)
        # Fall back to blank layout if not available
        try:
            slide = self.prs.slides.add_slide(self.slide_layouts[3])
        except IndexError:
            slide = self.prs.slides.add_slide(self.slide_layouts[5])
            logger.warning("Two-column layout not found, using blank layout")

        slide.shapes.title.text = title

        # Left column
        if len(slide.placeholders) > 1:
            left_frame = slide.placeholders[1].text_frame
            left_frame.clear()
            for item in left_content:
                p = left_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

        # Right column
        if len(slide.placeholders) > 2:
            right_frame = slide.placeholders[2].text_frame
            right_frame.clear()
            for item in right_content:
                p = right_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

        logger.info(f"Added two-column slide: {title}")
        self.slides_created += 1
        return self

    def add_table_slide(
        self,
        title: str,
        headers: List[str],
        rows: List[List[str]],
        header_color: tuple = (68, 114, 196)
    ) -> "PPTGenerator":
        """
        Add a table slide.

        Args:
            title: Slide title
            headers: List of column headers
            rows: List of data rows (each row is a list of cell values)
            header_color: RGB tuple for header background (default: blue)

        Returns:
            Self for method chaining
        """
        slide = self.prs.slides.add_slide(self.slide_layouts[5])  # Blank layout

        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title

        # Calculate table dimensions
        rows_count = len(rows) + 1  # +1 for header
        cols_count = len(headers)

        # Position and size
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(0.5 * rows_count)

        # Add table
        table = slide.shapes.add_table(
            rows_count, cols_count, left, top, width, height
        ).table

        # Set headers
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header

            # Style header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*header_color)

            # Style text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(14)
                if paragraph.font.color:
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)

        # Fill data
        for row_idx, row_data in enumerate(rows, start=1):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_data)
                cell.text_frame.paragraphs[0].font.size = Pt(12)

        logger.info(f"Added table slide: {title} ({len(rows)}x{len(headers)})")
        self.slides_created += 1
        return self

    def add_image_slide(
        self,
        title: str,
        image_path: str,
        width: float = 6.0,
        caption: Optional[str] = None
    ) -> "PPTGenerator":
        """
        Add an image slide.

        Args:
            title: Slide title
            image_path: Path to image file
            width: Image width in inches (default: 6.0)
            caption: Optional image caption

        Returns:
            Self for method chaining
        """
        slide = self.prs.slides.add_slide(self.slide_layouts[5])
        slide.shapes.title.text = title

        # Add image (centered)
        left = Inches((10 - width) / 2)
        top = Inches(2)

        try:
            pic = slide.shapes.add_picture(
                image_path, left, top, width=Inches(width)
            )
            logger.info(f"Added image slide: {title} ({image_path})")

            # Add caption if provided
            if caption:
                txBox = slide.shapes.add_textbox(
                    Inches(1), Inches(6.5), Inches(8), Inches(0.5)
                )
                tf = txBox.text_frame
                tf.text = caption
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                tf.paragraphs[0].font.size = Pt(12)
                tf.paragraphs[0].font.italic = True

        except FileNotFoundError:
            logger.error(f"Image not found: {image_path}")

        self.slides_created += 1
        return self

    def add_chart_slide(
        self,
        title: str,
        chart_type: str,
        data: Dict[str, Any]
    ) -> "PPTGenerator":
        """
        Add a chart slide.

        Args:
            title: Slide title
            chart_type: Type of chart ('line', 'bar', 'column', 'pie')
            data: Chart data dictionary with 'categories' and 'series'
                Example: {
                    'categories': ['Q1', 'Q2', 'Q3', 'Q4'],
                    'series': [
                        {'name': '2024', 'values': [100, 120, 150, 180]},
                        {'name': '2025', 'values': [120, 150, 180, 220]}
                    ]
                }

        Returns:
            Self for method chaining
        """
        slide = self.prs.slides.add_slide(self.slide_layouts[5])
        slide.shapes.title.text = title

        # Map chart type to enum
        chart_type_map = {
            'line': XL_CHART_TYPE.LINE,
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'pie': XL_CHART_TYPE.PIE
        }

        chart_enum = chart_type_map.get(chart_type, XL_CHART_TYPE.LINE)

        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = data.get('categories', [])

        for series in data.get('series', []):
            chart_data.add_series(
                series.get('name', 'Series'),
                series.get('values', [])
            )

        # Add chart
        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)
        chart = slide.shapes.add_chart(
            chart_enum, x, y, cx, cy, chart_data
        ).chart

        # Style chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM

        logger.info(f"Added {chart_type} chart slide: {title}")
        self.slides_created += 1
        return self

    def save(self, output_path: str) -> Dict[str, Any]:
        """
        Save the presentation to a file.

        Args:
            output_path: Path where to save the presentation

        Returns:
            Dictionary with success status and file info
        """
        try:
            # Create output directory if needed
            output_file = Path(output_path)
            output_file.parent.mkdir(parents=True, exist_ok=True)

            # Save presentation
            self.prs.save(str(output_file))

            file_size = output_file.stat().st_size

            logger.info(f"Presentation saved: {output_path}")
            logger.info(f"Slides: {self.slides_created}, Size: {file_size:,} bytes")

            return {
                "success": True,
                "file_path": str(output_file.absolute()),
                "slide_count": self.slides_created,
                "size_bytes": file_size
            }

        except Exception as e:
            logger.error(f"Failed to save presentation: {e}")
            return {
                "success": False,
                "error": str(e)
            }


def create_presentation(
    title: str,
    slides: List[Dict[str, Any]],
    output_path: str,
    subtitle: Optional[str] = None,
    template_path: Optional[str] = None
) -> Dict[str, Any]:
    """
    Create a presentation from structured data.

    Args:
        title: Presentation title
        slides: List of slide dictionaries
        output_path: Where to save the presentation
        subtitle: Optional subtitle
        template_path: Optional template file

    Returns:
        Result dictionary with success status
    """
    generator = PPTGenerator(template_path)

    # Add title slide
    generator.add_title_slide(title, subtitle)

    # Process slides
    for slide in slides:
        slide_type = slide.get('type', 'content')

        if slide_type == 'content':
            generator.add_content_slide(
                slide.get('title', ''),
                slide.get('content', [])
            )

        elif slide_type == 'two_column':
            generator.add_two_column_slide(
                slide.get('title', ''),
                slide.get('left', []),
                slide.get('right', [])
            )

        elif slide_type == 'table':
            generator.add_table_slide(
                slide.get('title', ''),
                slide.get('headers', []),
                slide.get('rows', [])
            )

        elif slide_type == 'image':
            generator.add_image_slide(
                slide.get('title', ''),
                slide.get('image_path', ''),
                slide.get('width', 6.0),
                slide.get('caption')
            )

        elif slide_type == 'chart':
            generator.add_chart_slide(
                slide.get('title', ''),
                slide.get('chart_type', 'line'),
                slide.get('data', {})
            )

    return generator.save(output_path)


def create_from_template(
    template_type: str,
    data: Dict[str, Any],
    output_path: str
) -> Dict[str, Any]:
    """
    Create presentation from predefined template.

    Args:
        template_type: Template name ('business-report', 'product-launch', etc.)
        data: Template data
        output_path: Where to save the presentation

    Returns:
        Result dictionary with success status
    """
    # Import template registry
    import templates

    try:
        template = templates.TemplateRegistry.get_template(template_type)
        slides = template.render(data)

        return create_presentation(
            title=data.get('title', 'Presentation'),
            slides=slides,
            output_path=output_path,
            subtitle=data.get('subtitle')
        )

    except KeyError:
        return {
            "success": False,
            "error": f"Template not found: {template_type}"
        }
