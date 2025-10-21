"""
HTML to PowerPoint Converter

This module provides HTML to PPT conversion functionality using BeautifulSoup
and python-pptx. It parses HTML structure and converts it to PowerPoint slides.
"""

import logging
from pathlib import Path
from typing import Optional, Dict, Any

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class HTMLtoPPTConverter:
    """
    HTML to PowerPoint converter.

    Converts HTML content to PowerPoint presentations by parsing HTML structure
    and mapping it to appropriate slide layouts.

    HTML Structure:
        - Each <section> becomes one slide
        - <h1> becomes slide title
        - <p> becomes paragraph text
        - <ul>/<ol> becomes bullet/numbered lists
        - <table> becomes PowerPoint table
        - <img> becomes image (requires src attribute)
    """

    def __init__(self, template_path: Optional[str] = None):
        """
        Initialize the HTML to PPT converter.

        Args:
            template_path: Path to PowerPoint template file (.pptx).
                          If None, creates a blank presentation.
        """
        if template_path:
            self.prs = Presentation(template_path)
            logger.info(f"Loaded template: {template_path}")
        else:
            self.prs = Presentation()
            logger.info("Created new blank presentation")

        self.slide_layouts = self.prs.slide_layouts
        self.slides_created = 0

    def convert(self, html_content: str) -> Presentation:
        """
        Convert HTML content to PowerPoint presentation.

        Args:
            html_content: HTML string to convert

        Returns:
            Presentation object
        """
        soup = BeautifulSoup(html_content, 'lxml')

        # Find all sections (each section = one slide)
        sections = soup.find_all('section')

        if not sections:
            logger.warning("No <section> tags found, treating entire HTML as one slide")
            self._process_section(soup)
        else:
            for section in sections:
                self._process_section(section)

        logger.info(f"Converted HTML to {self.slides_created} slides")
        return self.prs

    def _process_section(self, section):
        """Process a section element and create a slide."""
        # Determine slide type based on content
        title_elem = section.find(['h1', 'h2'])
        has_table = section.find('table') is not None
        has_image = section.find('img') is not None

        if has_table:
            self._create_table_slide(section, title_elem)
        elif has_image:
            self._create_image_slide(section, title_elem)
        else:
            self._create_content_slide(section, title_elem)

    def _create_content_slide(self, section, title_elem):
        """Create a content slide with title and bullet points."""
        slide = self.prs.slides.add_slide(self.slide_layouts[1])  # Content layout

        # Set title
        if title_elem:
            slide.shapes.title.text = title_elem.get_text(strip=True)

        # Get content placeholder
        if len(slide.placeholders) > 1:
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()

            # Process paragraphs
            for p in section.find_all('p'):
                paragraph = text_frame.add_paragraph()
                paragraph.text = p.get_text(strip=True)
                paragraph.level = 0
                paragraph.font.size = Pt(18)

            # Process lists
            for ul in section.find_all(['ul', 'ol']):
                for li in ul.find_all('li'):
                    paragraph = text_frame.add_paragraph()
                    paragraph.text = li.get_text(strip=True)
                    paragraph.level = 0
                    paragraph.font.size = Pt(16)

        self.slides_created += 1
        logger.info(f"Created content slide: {title_elem.get_text(strip=True) if title_elem else 'Untitled'}")

    def _create_table_slide(self, section, title_elem):
        """Create a slide with a table."""
        slide = self.prs.slides.add_slide(self.slide_layouts[5])  # Blank layout

        # Set title
        if title_elem:
            title_shape = slide.shapes.title
            title_shape.text = title_elem.get_text(strip=True)

        # Find table
        table_elem = section.find('table')
        if not table_elem:
            return

        # Extract table data
        headers = []
        rows_data = []

        # Get headers from <th> or first <tr>
        thead = table_elem.find('thead')
        if thead:
            header_row = thead.find('tr')
            if header_row:
                headers = [th.get_text(strip=True) for th in header_row.find_all(['th', 'td'])]

        # If no thead, try first row
        if not headers:
            first_row = table_elem.find('tr')
            if first_row:
                cells = first_row.find_all(['th', 'td'])
                if cells and cells[0].name == 'th':
                    headers = [cell.get_text(strip=True) for cell in cells]

        # Get body rows
        tbody = table_elem.find('tbody')
        rows_to_process = tbody.find_all('tr') if tbody else table_elem.find_all('tr')

        # Skip first row if it was used as headers
        start_index = 1 if headers and not thead else 0

        for tr in rows_to_process[start_index:]:
            cells = tr.find_all(['td', 'th'])
            row_data = [cell.get_text(strip=True) for cell in cells]
            if row_data:
                rows_data.append(row_data)

        # Create PowerPoint table
        if headers or rows_data:
            rows_count = len(rows_data) + (1 if headers else 0)
            cols_count = len(headers) if headers else len(rows_data[0]) if rows_data else 0

            if cols_count > 0:
                # Position and size
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9)
                height = Inches(0.5 * rows_count)

                # Add table
                table = slide.shapes.add_table(
                    rows_count, cols_count, left, top, width, height
                ).table

                # Set headers
                if headers:
                    for col_idx, header in enumerate(headers):
                        cell = table.cell(0, col_idx)
                        cell.text = header

                        # Style header
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)

                        # Style text
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.bold = True
                            paragraph.font.size = Pt(14)
                            if paragraph.font.color:
                                paragraph.font.color.rgb = RGBColor(255, 255, 255)

                # Fill data
                start_row = 1 if headers else 0
                for row_idx, row_data in enumerate(rows_data):
                    for col_idx, cell_data in enumerate(row_data[:cols_count]):
                        cell = table.cell(start_row + row_idx, col_idx)
                        cell.text = str(cell_data)
                        cell.text_frame.paragraphs[0].font.size = Pt(12)

        self.slides_created += 1
        logger.info(f"Created table slide: {title_elem.get_text(strip=True) if title_elem else 'Untitled'}")

    def _create_image_slide(self, section, title_elem):
        """Create a slide with an image."""
        slide = self.prs.slides.add_slide(self.slide_layouts[5])  # Blank layout

        # Set title
        if title_elem:
            slide.shapes.title.text = title_elem.get_text(strip=True)

        # Find image
        img = section.find('img')
        if img and img.get('src'):
            image_path = img.get('src')

            # Try to add image
            try:
                # Center image
                left = Inches(1.5)
                top = Inches(2)
                width = Inches(7)

                pic = slide.shapes.add_picture(
                    image_path, left, top, width=width
                )

                # Add caption if alt text exists
                if img.get('alt'):
                    txBox = slide.shapes.add_textbox(
                        Inches(1), Inches(6.5), Inches(8), Inches(0.5)
                    )
                    tf = txBox.text_frame
                    tf.text = img.get('alt')
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    tf.paragraphs[0].font.size = Pt(12)
                    tf.paragraphs[0].font.italic = True

                logger.info(f"Added image: {image_path}")

            except FileNotFoundError:
                logger.error(f"Image not found: {image_path}")
            except Exception as e:
                logger.error(f"Failed to add image {image_path}: {e}")

        self.slides_created += 1
        logger.info(f"Created image slide: {title_elem.get_text(strip=True) if title_elem else 'Untitled'}")

    def save(self, output_path: str) -> Dict[str, Any]:
        """
        Save the presentation to a file.

        Args:
            output_path: Path where to save the presentation

        Returns:
            Dictionary with success status and file info
        """
        try:
            # Create output directory if needed
            output_file = Path(output_path)
            output_file.parent.mkdir(parents=True, exist_ok=True)

            # Save presentation
            self.prs.save(str(output_file))

            file_size = output_file.stat().st_size

            logger.info(f"Presentation saved: {output_path}")
            logger.info(f"Slides: {self.slides_created}, Size: {file_size:,} bytes")

            return {
                "success": True,
                "file_path": str(output_file.absolute()),
                "slide_count": self.slides_created,
                "size_bytes": file_size
            }

        except Exception as e:
            logger.error(f"Failed to save presentation: {e}")
            return {
                "success": False,
                "error": str(e)
            }


def convert_html_to_ppt(
    html_content: str,
    output_path: str,
    template_path: Optional[str] = None
) -> Dict[str, Any]:
    """
    Convert HTML to PowerPoint presentation.

    Args:
        html_content: HTML string to convert
        output_path: Where to save the presentation
        template_path: Optional template file

    Returns:
        Result dictionary with success status
    """
    converter = HTMLtoPPTConverter(template_path)
    converter.convert(html_content)
    return converter.save(output_path)
