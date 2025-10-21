#!/usr/bin/env python3
"""
将截图组合成PPT - 每页一张图片
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

# 路径配置
screenshot_dir = Path(__file__).parent / "screenshots_paginated"
output_ppt = Path(__file__).parent / "云南过桥米线营销方案_截图版.pptx"

# PPT尺寸 (16:9宽屏)
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)

def create_ppt_from_screenshots():
    """将截图组合成PPT"""

    # 获取所有截图
    screenshots = sorted(screenshot_dir.glob("slide_*.png"))

    if not screenshots:
        print("❌ 未找到截图文件")
        return

    print(f"📸 找到 {len(screenshots)} 张截图")
    print()

    # 创建PPT
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 添加空白版式
    blank_slide_layout = prs.slide_layouts[6]  # 空白版式

    for screenshot in screenshots:
        print(f"  ➕ 添加: {screenshot.name}")

        # 创建新幻灯片
        slide = prs.slides.add_slide(blank_slide_layout)

        # 添加图片,填满整个幻灯片
        slide.shapes.add_picture(
            str(screenshot),
            left=0,
            top=0,
            width=SLIDE_WIDTH,
            height=SLIDE_HEIGHT
        )

    # 保存PPT
    prs.save(str(output_ppt))

    print()
    print(f"✅ PPT已生成: {output_ppt}")
    print(f"📊 共 {len(screenshots)} 页")
    print()
    print("🎉 完成!现在您可以打开PPT查看效果")

if __name__ == "__main__":
    create_ppt_from_screenshots()
